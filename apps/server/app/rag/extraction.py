"""Turn uploaded bytes into indexable text.

Parsing is CPU-bound and synchronous, so callers must go through
`extract_text`, which offloads to a worker thread rather than blocking the
event loop.
"""

import io
import logging

from fastapi.concurrency import run_in_threadpool

logger = logging.getLogger(__name__)

TEXT_MIMES = {
    "text/plain",
    "text/markdown",
    "text/csv",
    "text/html",
    "application/json",
    "application/xml",
}

PDF_MIMES = {"application/pdf"}

DOCX_MIMES = {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}
PPTX_MIMES = {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}
XLSX_MIMES = {"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"}

IMAGE_MIMES = {"image/png", "image/jpeg", "image/jpg", "image/webp", "image/tiff", "image/bmp"}

# A text layer this thin means a scanned document, not an empty one.
_MIN_TEXT_CHARS = 32

# Render free tier has one shared, low-CPU process running the API inline
# with ingestion. OCR-ing an unbounded page count can peg that process for
# minutes and stall concurrent chat/search requests, so scanned PDFs beyond
# this are reported as needing OCR rather than run through it here.
_MAX_OCR_PAGES = 20


class Extraction:
    """Result of parsing an upload."""

    def __init__(self, content: str | None, doc_type: str, status: str) -> None:
        self.content = content
        self.doc_type = doc_type
        self.status = status


def _extract_pdf(raw: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(raw))
    pages = []
    for page in reader.pages:
        try:
            pages.append(page.extract_text() or "")
        except Exception:
            # One malformed page should not lose the rest of the document.
            logger.warning("failed to extract a PDF page", exc_info=True)
    return "\n\n".join(p.strip() for p in pages if p.strip()).strip()


def _decode_text(raw: bytes) -> str:
    return raw.decode("utf-8", errors="replace").strip()


def _extract_docx(raw: bytes) -> str:
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(raw))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts).strip()


def _extract_pptx(raw: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    slides = []
    for slide in prs.slides:
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if texts:
            slides.append("\n".join(texts))
    return "\n\n".join(slides).strip()


def _extract_xlsx(raw: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"# {ws.title}\n" + "\n".join(rows))
    return "\n\n".join(sheets).strip()


def _ocr_image(raw: bytes) -> str:
    import pytesseract
    from PIL import Image

    return pytesseract.image_to_string(Image.open(io.BytesIO(raw))).strip()


def _ocr_pdf(raw: bytes) -> str:
    import pytesseract
    from pdf2image import convert_from_bytes
    from pypdf import PdfReader

    page_count = len(PdfReader(io.BytesIO(raw)).pages)
    if page_count > _MAX_OCR_PAGES:
        raise ValueError(
            f"PDF has {page_count} pages, over the {_MAX_OCR_PAGES}-page OCR limit"
        )

    pages = convert_from_bytes(raw)
    texts = [pytesseract.image_to_string(page) for page in pages]
    return "\n\n".join(t.strip() for t in texts if t.strip()).strip()


async def extract_text(raw: bytes, mime: str, filename: str = "") -> Extraction:
    """Map raw upload bytes to (content, document type, ingest status)."""
    mime = (mime or "").split(";")[0].strip().lower()
    suffix = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if mime in PDF_MIMES or suffix == "pdf":
        try:
            content = await run_in_threadpool(_extract_pdf, raw)
        except Exception:
            logger.exception("PDF extraction failed for %s", filename)
            return Extraction(None, "pdf", "failed")
        if len(content) >= _MIN_TEXT_CHARS:
            return Extraction(content, "pdf", "pending")
        # Almost certainly a scan. Try OCR; fall back to needs_ocr if the
        # runtime has no tesseract binary (e.g. local dev without it installed).
        try:
            ocr_content = await run_in_threadpool(_ocr_pdf, raw)
        except Exception:
            logger.warning("OCR unavailable or failed for %s", filename, exc_info=True)
            return Extraction(None, "pdf", "needs_ocr")
        if len(ocr_content) < _MIN_TEXT_CHARS:
            return Extraction(None, "pdf", "needs_ocr")
        return Extraction(ocr_content, "pdf", "pending")

    if mime in TEXT_MIMES or suffix in {"txt", "md", "csv", "json", "xml", "html"}:
        content = await run_in_threadpool(_decode_text, raw)
        if not content:
            return Extraction(None, "note", "skipped_empty")
        return Extraction(content, "note", "pending")

    if mime in DOCX_MIMES or suffix == "docx":
        try:
            content = await run_in_threadpool(_extract_docx, raw)
        except Exception:
            logger.exception("DOCX extraction failed for %s", filename)
            return Extraction(None, "docx", "failed")
        if not content:
            return Extraction(None, "docx", "skipped_empty")
        return Extraction(content, "docx", "pending")

    if mime in PPTX_MIMES or suffix == "pptx":
        try:
            content = await run_in_threadpool(_extract_pptx, raw)
        except Exception:
            logger.exception("PPTX extraction failed for %s", filename)
            return Extraction(None, "pptx", "failed")
        if not content:
            return Extraction(None, "pptx", "skipped_empty")
        return Extraction(content, "pptx", "pending")

    if mime in XLSX_MIMES or suffix == "xlsx":
        try:
            content = await run_in_threadpool(_extract_xlsx, raw)
        except Exception:
            logger.exception("XLSX extraction failed for %s", filename)
            return Extraction(None, "xlsx", "failed")
        if not content:
            return Extraction(None, "xlsx", "skipped_empty")
        return Extraction(content, "xlsx", "pending")

    if mime in IMAGE_MIMES or suffix in {"png", "jpg", "jpeg", "webp", "tiff", "bmp"}:
        try:
            content = await run_in_threadpool(_ocr_image, raw)
        except Exception:
            logger.warning("OCR unavailable or failed for %s", filename, exc_info=True)
            return Extraction(None, "image", "needs_ocr")
        if len(content) < _MIN_TEXT_CHARS:
            return Extraction(None, "image", "needs_ocr")
        return Extraction(content, "image", "pending")

    # Stored but not indexable — archives and anything else unrecognized.
    return Extraction(None, "file", "unsupported")
